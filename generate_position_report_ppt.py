#!/usr/bin/env python3
"""
Generate Strategy A-G Position Building Report PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import json
import os
from datetime import datetime

# === Colors ===
BG_DARK = RGBColor(0x1a, 0x1a, 0x2e)
BG_CARD = RGBColor(0x16, 0x21, 0x3e)
ACCENT_BLUE = RGBColor(0x00, 0xd2, 0xff)
ACCENT_GREEN = RGBColor(0x00, 0xe6, 0x76)
ACCENT_RED = RGBColor(0xff, 0x41, 0x57)
ACCENT_GOLD = RGBColor(0xff, 0xd7, 0x00)
ACCENT_PURPLE = RGBColor(0xbb, 0x86, 0xfc)
TEXT_WHITE = RGBColor(0xff, 0xff, 0xff)
TEXT_GRAY = RGBColor(0xaa, 0xaa, 0xcc)
TEXT_LIGHT = RGBColor(0xe0, 0xe0, 0xf0)

DATA_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'data')

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=14, color=TEXT_WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='Arial'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_card(slide, left, top, width, height, color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_metric_card(slide, left, top, label, value, color=ACCENT_BLUE):
    card = add_card(slide, left, top, 2.0, 0.9)
    add_text_box(slide, left+0.1, top+0.05, 1.8, 0.35, label, 10, TEXT_GRAY, font_name='Microsoft YaHei')
    add_text_box(slide, left+0.1, top+0.35, 1.8, 0.5, value, 20, color, True, font_name='Arial')

def create_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    # Title
    add_text_box(slide, 0.5, 1.0, 9.0, 1.0, '量化策略建仓报告', 40, ACCENT_BLUE, True, PP_ALIGN.CENTER, 'Microsoft YaHei')
    add_text_box(slide, 0.5, 1.8, 9.0, 0.6, 'Strategy A ~ G Portfolio Position Report', 18, TEXT_GRAY, False, PP_ALIGN.CENTER)

    # Date
    add_text_box(slide, 0.5, 2.8, 9.0, 0.5, f'Signal Date: 2026-03-20  |  Report: {datetime.now().strftime("%Y-%m-%d")}', 14, TEXT_LIGHT, False, PP_ALIGN.CENTER)

    # Strategy overview cards
    strategies = [
        ('A', '纯债ETF轮动', '4.1%', '-2.1%', ACCENT_GREEN),
        ('B', '4因子ETF动量', '15.6%', '-18.3%', ACCENT_BLUE),
        ('C', '港股Top1轮动', '17.0%', '-23.3%', ACCENT_PURPLE),
        ('D', 'CSI300低波动', '13.9%', '-8.5%', ACCENT_GOLD),
        ('E', '价值反转', '22.0%', '-17.6%', ACCENT_RED),
        ('F', '美股QDII动量', '13.0%', '-18.9%', ACCENT_BLUE),
        ('G', 'CSI500低波价值', '13.5%', '-11.9%', ACCENT_GREEN),
    ]

    y = 3.6
    for i, (code, name, cagr, mdd, color) in enumerate(strategies):
        x = 0.3 + (i % 4) * 2.4
        if i == 4:
            y = 4.65
            x = 0.3 + ((i-4) % 4) * 2.4
        if i > 4:
            x = 0.3 + ((i-4) % 4) * 2.4

        card = add_card(slide, x, y, 2.2, 0.85)
        add_text_box(slide, x+0.1, y+0.05, 2.0, 0.3, f'{code}: {name}', 11, color, True, font_name='Microsoft YaHei')
        add_text_box(slide, x+0.1, y+0.4, 1.0, 0.3, f'CAGR {cagr}', 10, ACCENT_GREEN, font_name='Arial')
        add_text_box(slide, x+1.1, y+0.4, 1.0, 0.3, f'MDD {mdd}', 10, ACCENT_RED, font_name='Arial')

def create_strategy_a_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, 0.5, 0.3, 9.0, 0.6, 'A  纯债ETF动量轮动', 28, ACCENT_GREEN, True, PP_ALIGN.LEFT, 'Microsoft YaHei')
    add_text_box(slide, 0.5, 0.8, 9.0, 0.4, '四只债券ETF按8周动量轮动 | 15年回测 | Calmar 1.97', 12, TEXT_GRAY, font_name='Microsoft YaHei')

    # Metrics
    add_metric_card(slide, 0.5, 1.4, 'CAGR', '4.1%', ACCENT_GREEN)
    add_metric_card(slide, 2.7, 1.4, 'Max DD', '-2.1%', ACCENT_RED)
    add_metric_card(slide, 4.9, 1.4, 'Sharpe', '1.045', ACCENT_BLUE)
    add_metric_card(slide, 7.1, 1.4, 'Calmar', '1.971', ACCENT_GOLD)

    # Holdings
    add_card(slide, 0.5, 2.6, 9.0, 2.8)
    add_text_box(slide, 0.7, 2.7, 8.5, 0.4, '当前信号 (3/20)', 16, ACCENT_BLUE, True, font_name='Microsoft YaHei')

    holdings = [
        ('511260', '十年国债ETF', '50%', '动量#1'),
        ('511220', '城投债ETF', '30%', '动量#2'),
        ('511030', '信用债ETF', '20%', '动量#3'),
        ('511360', '短债ETF', '0%', '动量#4'),
    ]
    for i, (code, name, weight, rank) in enumerate(holdings):
        y = 3.2 + i * 0.5
        color = ACCENT_GREEN if weight != '0%' else TEXT_GRAY
        add_text_box(slide, 0.7, y, 1.5, 0.4, code, 13, color, True, font_name='Arial')
        add_text_box(slide, 2.2, y, 2.5, 0.4, name, 13, TEXT_WHITE, font_name='Microsoft YaHei')
        add_text_box(slide, 5.0, y, 1.5, 0.4, weight, 16, color, True, PP_ALIGN.CENTER, font_name='Arial')
        add_text_box(slide, 6.8, y, 2.0, 0.4, rank, 11, TEXT_GRAY, font_name='Arial')

def create_strategy_b_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, 0.5, 0.3, 9.0, 0.6, 'B  4因子ETF动量轮动', 28, ACCENT_BLUE, True, PP_ALIGN.LEFT, 'Microsoft YaHei')
    add_text_box(slide, 0.5, 0.8, 9.0, 0.4, '20万实盘 | 4因子池Top3集中 | Vol缩放 | 周度调仓', 12, TEXT_GRAY, font_name='Microsoft YaHei')

    add_metric_card(slide, 0.5, 1.4, 'CAGR', '15.6%', ACCENT_GREEN)
    add_metric_card(slide, 2.7, 1.4, 'Max DD', '-18.3%', ACCENT_RED)
    add_metric_card(slide, 4.9, 1.4, 'Sharpe', '1.128', ACCENT_BLUE)
    add_metric_card(slide, 7.1, 1.4, '投入资金', '20万', ACCENT_GOLD)

    add_card(slide, 0.5, 2.6, 9.0, 3.2)
    add_text_box(slide, 0.7, 2.7, 8.5, 0.4, '建仓计划 (3/24 周一开盘)', 16, ACCENT_BLUE, True, font_name='Microsoft YaHei')
    add_text_box(slide, 0.7, 3.1, 8.5, 0.3, '市场状态: 过渡 | 总仓位: 91% (Regime 70% x Vol缩放 130%)', 11, TEXT_GRAY, font_name='Microsoft YaHei')

    holdings = [
        ('515080', '红利低波ETF', '45.5%', '91,000元', '#1 +2.1%'),
        ('159201', '自由现金流ETF', '27.3%', '54,600元', '#2 -2.0%'),
        ('512040', '国信价值ETF', '18.2%', '36,400元', '#3 -3.0%'),
        ('159967', '创业板成长ETF', '0%', '0元', '#4 -3.7%'),
        ('现金', '货币基金', '9.0%', '18,000元', ''),
    ]
    for i, (code, name, weight, amount, rank) in enumerate(holdings):
        y = 3.5 + i * 0.5
        color = ACCENT_GREEN if weight not in ('0%', '9.0%') else TEXT_GRAY
        if weight == '9.0%': color = ACCENT_GOLD
        add_text_box(slide, 0.7, y, 1.2, 0.4, code, 13, color, True, font_name='Arial')
        add_text_box(slide, 1.9, y, 2.2, 0.4, name, 13, TEXT_WHITE, font_name='Microsoft YaHei')
        add_text_box(slide, 4.3, y, 1.0, 0.4, weight, 14, color, True, PP_ALIGN.CENTER, font_name='Arial')
        add_text_box(slide, 5.5, y, 1.8, 0.4, amount, 13, ACCENT_GOLD, True, PP_ALIGN.RIGHT, font_name='Arial')
        add_text_box(slide, 7.5, y, 1.8, 0.4, rank, 11, TEXT_GRAY, font_name='Arial')

def create_strategy_c_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, 0.5, 0.3, 9.0, 0.6, 'C  港股ETF纯动量轮动', 28, ACCENT_PURPLE, True, PP_ALIGN.LEFT, 'Microsoft YaHei')
    add_text_box(slide, 0.5, 0.8, 9.0, 0.4, '3只港股ETF按4周动量Top1全仓 | 回测CAGR 17%', 12, TEXT_GRAY, font_name='Microsoft YaHei')

    add_metric_card(slide, 0.5, 1.4, 'CAGR', '17.0%', ACCENT_GREEN)
    add_metric_card(slide, 2.7, 1.4, 'Max DD', '-23.3%', ACCENT_RED)
    add_metric_card(slide, 4.9, 1.4, 'Sharpe', '~0.85', ACCENT_BLUE)
    add_metric_card(slide, 7.1, 1.4, '持仓', 'Top1全仓', ACCENT_GOLD)

    add_card(slide, 0.5, 2.6, 9.0, 2.5)
    add_text_box(slide, 0.7, 2.7, 8.5, 0.4, '当前信号 (3/20)', 16, ACCENT_PURPLE, True, font_name='Microsoft YaHei')

    holdings = [
        ('159726', '恒生高股息ETF', '100%', '-1.6%', True),
        ('513180', '恒生科技ETF', '0%', '-11.0%', False),
        ('513060', '恒生医疗ETF', '0%', '-11.6%', False),
    ]
    for i, (code, name, weight, mom, selected) in enumerate(holdings):
        y = 3.2 + i * 0.55
        color = ACCENT_GREEN if selected else TEXT_GRAY
        marker = '>>' if selected else '  '
        add_text_box(slide, 0.7, y, 0.3, 0.4, marker, 14, color, True, font_name='Arial')
        add_text_box(slide, 1.1, y, 1.5, 0.4, code, 14, color, True, font_name='Arial')
        add_text_box(slide, 2.6, y, 2.5, 0.4, name, 14, TEXT_WHITE if selected else TEXT_GRAY, font_name='Microsoft YaHei')
        add_text_box(slide, 5.5, y, 1.2, 0.4, weight, 16, color, True, PP_ALIGN.CENTER, font_name='Arial')
        add_text_box(slide, 7.0, y, 2.0, 0.4, f'4w: {mom}', 12, TEXT_GRAY, font_name='Arial')

def create_strategy_d_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, 0.5, 0.3, 9.0, 0.6, 'D  CSI300 低波动', 28, ACCENT_GOLD, True, PP_ALIGN.LEFT, 'Microsoft YaHei')
    add_text_box(slide, 0.5, 0.8, 9.0, 0.4, '沪深300真实成分股 | 20周最低波动率Top10 | 无幸存者偏差', 12, TEXT_GRAY, font_name='Microsoft YaHei')

    add_metric_card(slide, 0.5, 1.3, 'CAGR', '13.9%', ACCENT_GREEN)
    add_metric_card(slide, 2.7, 1.3, 'Max DD', '-8.5%', ACCENT_RED)
    add_metric_card(slide, 4.9, 1.3, 'Sharpe', '0.955', ACCENT_BLUE)
    add_metric_card(slide, 7.1, 1.3, 'Calmar', '1.636', ACCENT_GOLD)

    add_card(slide, 0.5, 2.5, 9.0, 3.5)
    add_text_box(slide, 0.7, 2.55, 8.5, 0.35, '当前Top10持仓 (3/20) 等权10%', 14, ACCENT_GOLD, True, font_name='Microsoft YaHei')

    # Header
    add_text_box(slide, 0.7, 2.9, 0.5, 0.3, '#', 9, TEXT_GRAY, True, font_name='Arial')
    add_text_box(slide, 1.1, 2.9, 1.5, 0.3, '代码', 9, TEXT_GRAY, True, font_name='Microsoft YaHei')
    add_text_box(slide, 2.5, 2.9, 2.0, 0.3, '名称', 9, TEXT_GRAY, True, font_name='Microsoft YaHei')
    add_text_box(slide, 4.5, 2.9, 1.5, 0.3, '行业', 9, TEXT_GRAY, True, font_name='Microsoft YaHei')
    add_text_box(slide, 6.2, 2.9, 1.0, 0.3, '波动率', 9, TEXT_GRAY, True, PP_ALIGN.CENTER, font_name='Microsoft YaHei')
    add_text_box(slide, 7.5, 2.9, 1.5, 0.3, '12w动量', 9, TEXT_GRAY, True, PP_ALIGN.CENTER, font_name='Microsoft YaHei')

    stocks = [
        (1, '000538', '云南白药', '医药', '7.8%', '+0.4%'),
        (2, '000651', '格力电器', '电气', '9.5%', '-4.8%'),
        (3, '601607', '上海医药', '零售', '9.8%', '-5.7%'),
        (4, '601059', '信达证券', '证券', '10.3%', '-4.8%'),
        (5, '601788', '光大证券', '证券', '10.5%', '-9.8%'),
        (6, '000001', '平安银行', '银行', '10.5%', '-7.3%'),
        (7, '600900', '长江电力', '电力', '10.8%', '-1.6%'),
        (8, '601878', '浙商证券', '证券', '10.8%', '-8.2%'),
        (9, '000166', '申万宏源', '证券', '11.0%', '-6.9%'),
        (10, '601816', '京沪高铁', '铁路', '11.1%', '-2.9%'),
    ]
    for i, (rank, code, name, industry, vol, mom) in enumerate(stocks):
        y = 3.15 + i * 0.28
        c = TEXT_WHITE if i < 3 else TEXT_LIGHT
        add_text_box(slide, 0.7, y, 0.4, 0.28, str(rank), 10, ACCENT_GOLD, True, font_name='Arial')
        add_text_box(slide, 1.1, y, 1.4, 0.28, code, 10, c, font_name='Arial')
        add_text_box(slide, 2.5, y, 2.0, 0.28, name, 10, c, font_name='Microsoft YaHei')
        add_text_box(slide, 4.5, y, 1.5, 0.28, industry, 10, TEXT_GRAY, font_name='Microsoft YaHei')
        add_text_box(slide, 6.2, y, 1.0, 0.28, vol, 10, ACCENT_GREEN, True, PP_ALIGN.CENTER, font_name='Arial')
        mom_color = ACCENT_GREEN if mom.startswith('+') else ACCENT_RED
        add_text_box(slide, 7.5, y, 1.5, 0.28, mom, 10, mom_color, False, PP_ALIGN.CENTER, font_name='Arial')

def create_strategy_e_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, 0.5, 0.3, 9.0, 0.6, 'E  龙头价值反转', 28, ACCENT_RED, True, PP_ALIGN.LEFT, 'Microsoft YaHei')
    add_text_box(slide, 0.5, 0.8, 9.0, 0.4, '28只A股龙头 | 价值(逆向)+质量 | Top4等权 | 月度换仓', 12, TEXT_GRAY, font_name='Microsoft YaHei')

    add_metric_card(slide, 0.5, 1.4, 'CAGR', '22.0%', ACCENT_GREEN)
    add_metric_card(slide, 2.7, 1.4, 'Max DD', '-17.6%', ACCENT_RED)
    add_metric_card(slide, 4.9, 1.4, 'Sharpe', '~1.1', ACCENT_BLUE)
    add_metric_card(slide, 7.1, 1.4, '偏差风险', '中等', ACCENT_GOLD)

    add_card(slide, 0.5, 2.6, 9.0, 2.8)
    add_text_box(slide, 0.7, 2.7, 8.5, 0.4, '当前Top4持仓 (3/20) 等权25%', 14, ACCENT_RED, True, font_name='Microsoft YaHei')

    stocks = [
        (1, '300760', '迈瑞医疗', '医疗器械', '-15.7%', '-8.8%'),
        (2, '002230', '科大讯飞', 'AI', '-23.0%', '-12.1%'),
        (3, '600276', '恒瑞医药', '创新药', '-11.0%', '-2.5%'),
        (4, '000063', '中兴通讯', '通信设备', '-14.4%', '-11.3%'),
    ]
    for i, (rank, code, name, sector, mom12, mom4) in enumerate(stocks):
        y = 3.2 + i * 0.55
        add_text_box(slide, 0.7, y, 0.4, 0.4, f'#{rank}', 14, ACCENT_RED, True, font_name='Arial')
        add_text_box(slide, 1.2, y, 1.2, 0.4, code, 14, TEXT_WHITE, True, font_name='Arial')
        add_text_box(slide, 2.5, y, 2.0, 0.4, name, 14, TEXT_WHITE, font_name='Microsoft YaHei')
        add_text_box(slide, 4.5, y, 1.5, 0.4, sector, 11, TEXT_GRAY, font_name='Microsoft YaHei')
        add_text_box(slide, 6.0, y, 1.5, 0.4, f'12w: {mom12}', 11, ACCENT_RED, font_name='Arial')
        add_text_box(slide, 7.6, y, 1.5, 0.4, f'4w: {mom4}', 11, ACCENT_RED, font_name='Arial')

def create_strategy_f_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, 0.5, 0.3, 9.0, 0.6, 'F  美股QDII动量轮动', 28, ACCENT_BLUE, True, PP_ALIGN.LEFT, 'Microsoft YaHei')
    add_text_box(slide, 0.5, 0.8, 9.0, 0.4, '纳指100+标普500 QDII ETF | 双均线+回撤锁定 | 避险切短债', 12, TEXT_GRAY, font_name='Microsoft YaHei')

    add_metric_card(slide, 0.5, 1.4, 'CAGR', '13.0%', ACCENT_GREEN)
    add_metric_card(slide, 2.7, 1.4, 'Max DD', '-18.9%', ACCENT_RED)
    add_metric_card(slide, 4.9, 1.4, 'Sharpe', '0.809', ACCENT_BLUE)
    add_metric_card(slide, 7.1, 1.4, 'Calmar', '0.689', ACCENT_GOLD)

    add_card(slide, 0.5, 2.6, 9.0, 2.5)
    add_text_box(slide, 0.7, 2.7, 8.5, 0.4, '当前信号', 16, ACCENT_BLUE, True, font_name='Microsoft YaHei')
    add_text_box(slide, 0.7, 3.15, 8.5, 0.3, '纳指100 8周MA > 20周MA → 进攻模式', 13, ACCENT_GREEN, font_name='Microsoft YaHei')

    holdings = [
        ('159941', '纳指100ETF', '50%', '动量#1 (QDII)'),
        ('513500', '标普500ETF', '30%', '动量#2 (QDII)'),
        ('511360', '短债ETF', '20%', '避险配置'),
    ]
    for i, (code, name, weight, note) in enumerate(holdings):
        y = 3.6 + i * 0.5
        color = ACCENT_GREEN if '动量' in note else ACCENT_GOLD
        add_text_box(slide, 0.7, y, 1.5, 0.4, code, 14, color, True, font_name='Arial')
        add_text_box(slide, 2.2, y, 2.5, 0.4, name, 14, TEXT_WHITE, font_name='Microsoft YaHei')
        add_text_box(slide, 5.0, y, 1.2, 0.4, weight, 16, color, True, PP_ALIGN.CENTER, font_name='Arial')
        add_text_box(slide, 6.5, y, 2.5, 0.4, note, 11, TEXT_GRAY, font_name='Arial')

def create_strategy_g_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, 0.5, 0.3, 9.0, 0.6, 'G  CSI500 低波价值', 28, ACCENT_GREEN, True, PP_ALIGN.LEFT, 'Microsoft YaHei')
    add_text_box(slide, 0.5, 0.8, 9.0, 0.4, '中证500真实成分股 | 50%低波+50%低PB | Top10 | 无幸存者偏差', 12, TEXT_GRAY, font_name='Microsoft YaHei')

    add_metric_card(slide, 0.5, 1.3, 'CAGR', '13.5%', ACCENT_GREEN)
    add_metric_card(slide, 2.7, 1.3, 'Max DD', '-11.9%', ACCENT_RED)
    add_metric_card(slide, 4.9, 1.3, 'Sharpe', '0.794', ACCENT_BLUE)
    add_metric_card(slide, 7.1, 1.3, 'Calmar', '1.132', ACCENT_GOLD)

    add_card(slide, 0.5, 2.5, 9.0, 3.5)
    add_text_box(slide, 0.7, 2.55, 8.5, 0.35, '当前Top10持仓 (3/20) 等权10%', 14, ACCENT_GREEN, True, font_name='Microsoft YaHei')

    # Header
    add_text_box(slide, 0.7, 2.9, 0.5, 0.3, '#', 9, TEXT_GRAY, True, font_name='Arial')
    add_text_box(slide, 1.1, 2.9, 1.5, 0.3, '代码', 9, TEXT_GRAY, True, font_name='Microsoft YaHei')
    add_text_box(slide, 2.5, 2.9, 2.0, 0.3, '名称', 9, TEXT_GRAY, True, font_name='Microsoft YaHei')
    add_text_box(slide, 4.5, 2.9, 1.5, 0.3, '行业', 9, TEXT_GRAY, True, font_name='Microsoft YaHei')
    add_text_box(slide, 6.0, 2.9, 0.8, 0.3, '波动率', 9, TEXT_GRAY, True, PP_ALIGN.CENTER, font_name='Microsoft YaHei')
    add_text_box(slide, 7.0, 2.9, 0.8, 0.3, 'PB', 9, TEXT_GRAY, True, PP_ALIGN.CENTER, font_name='Arial')
    add_text_box(slide, 8.0, 2.9, 1.0, 0.3, '评分', 9, TEXT_GRAY, True, PP_ALIGN.CENTER, font_name='Arial')

    stocks = [
        (1, '601997', '贵阳银行', '银行', '9.3%', '0.33', '1.19'),
        (2, '002958', '青农商行', '银行', '10.1%', '0.47', '1.15'),
        (3, '002966', '苏州银行', '银行', '9.9%', '0.73', '1.13'),
        (4, '600332', '白云山', '医药', '8.3%', '1.08', '1.11'),
        (5, '600109', '国金证券', '证券', '10.3%', '0.92', '1.09'),
        (6, '601128', '常熟银行', '银行', '11.5%', '0.79', '1.09'),
        (7, '002926', '华西证券', '证券', '11.1%', '0.92', '1.08'),
        (8, '601333', '广深铁路', '铁路', '12.0%', '0.77', '1.08'),
        (9, '601928', '凤凰传媒', '出版', '9.3%', '1.24', '1.07'),
        (10, '000750', '国海证券', '证券', '11.0%', '1.12', '1.05'),
    ]
    for i, (rank, code, name, industry, vol, pb, score) in enumerate(stocks):
        y = 3.15 + i * 0.28
        c = TEXT_WHITE if i < 3 else TEXT_LIGHT
        add_text_box(slide, 0.7, y, 0.4, 0.28, str(rank), 10, ACCENT_GREEN, True, font_name='Arial')
        add_text_box(slide, 1.1, y, 1.4, 0.28, code, 10, c, font_name='Arial')
        add_text_box(slide, 2.5, y, 2.0, 0.28, name, 10, c, font_name='Microsoft YaHei')
        add_text_box(slide, 4.5, y, 1.5, 0.28, industry, 10, TEXT_GRAY, font_name='Microsoft YaHei')
        add_text_box(slide, 6.0, y, 0.8, 0.28, vol, 10, ACCENT_GREEN, True, PP_ALIGN.CENTER, font_name='Arial')
        pb_color = ACCENT_GREEN if float(pb) < 1.0 else TEXT_LIGHT
        add_text_box(slide, 7.0, y, 0.8, 0.28, pb, 10, pb_color, True, PP_ALIGN.CENTER, font_name='Arial')
        add_text_box(slide, 8.0, y, 1.0, 0.28, score, 10, ACCENT_BLUE, False, PP_ALIGN.CENTER, font_name='Arial')

def create_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, 0.5, 0.3, 9.0, 0.6, '策略总览与建仓建议', 28, ACCENT_BLUE, True, PP_ALIGN.CENTER, 'Microsoft YaHei')

    # Summary table
    add_card(slide, 0.3, 1.2, 9.4, 3.8)

    # Header
    headers = ['策略', '类型', 'CAGR', 'MDD', 'Calmar', '信号概要']
    x_positions = [0.5, 1.3, 3.0, 4.0, 5.1, 6.2]
    widths = [0.8, 1.7, 1.0, 1.0, 1.0, 3.0]
    for h, x, w in zip(headers, x_positions, widths):
        add_text_box(slide, x, 1.3, w, 0.3, h, 10, ACCENT_BLUE, True, PP_ALIGN.CENTER, 'Microsoft YaHei')

    rows = [
        ('A', '纯债轮动', '4.1%', '-2.1%', '1.97', '十年国债50%+城投30%+信用20%'),
        ('B', '4因子ETF(实盘)', '15.6%', '-18.3%', '0.85', '红利低波45%+现金流27%+价值18%'),
        ('C', '港股Top1', '17.0%', '-23.3%', '0.73', '恒生高股息ETF 100%'),
        ('D', 'CSI300低波', '13.9%', '-8.5%', '1.64', '云南白药/格力/上海医药等10只'),
        ('E', '龙头价值反转', '22.0%', '-17.6%', '1.25', '迈瑞/讯飞/恒瑞/中兴 各25%'),
        ('F', '美股QDII', '13.0%', '-18.9%', '0.69', '纳指100+标普500+短债避险'),
        ('G', 'CSI500低波价值', '13.5%', '-11.9%', '1.13', '贵阳银行/青农商行/苏州银行等10只'),
    ]

    for i, (code, typ, cagr, mdd, calmar, signal) in enumerate(rows):
        y = 1.65 + i * 0.42
        colors = [ACCENT_GOLD, ACCENT_BLUE, ACCENT_PURPLE, ACCENT_GOLD, ACCENT_RED, ACCENT_BLUE, ACCENT_GREEN]
        add_text_box(slide, 0.5, y, 0.8, 0.35, code, 13, colors[i], True, PP_ALIGN.CENTER, 'Arial')
        add_text_box(slide, 1.3, y, 1.7, 0.35, typ, 10, TEXT_WHITE, False, PP_ALIGN.CENTER, 'Microsoft YaHei')
        add_text_box(slide, 3.0, y, 1.0, 0.35, cagr, 11, ACCENT_GREEN, True, PP_ALIGN.CENTER, 'Arial')
        add_text_box(slide, 4.0, y, 1.0, 0.35, mdd, 11, ACCENT_RED, False, PP_ALIGN.CENTER, 'Arial')
        add_text_box(slide, 5.1, y, 1.0, 0.35, calmar, 11, ACCENT_BLUE, True, PP_ALIGN.CENTER, 'Arial')
        add_text_box(slide, 6.2, y, 3.3, 0.35, signal, 9, TEXT_GRAY, False, PP_ALIGN.LEFT, 'Microsoft YaHei')

    # Notes
    add_text_box(slide, 0.5, 5.2, 9.0, 0.8,
        '建议: 确认各策略投入金额后，即可生成详细买入清单 (具体股数/手数/金额)\n'
        '策略B已确认20万实盘，本周一(3/24)开始建仓',
        11, TEXT_GRAY, False, PP_ALIGN.LEFT, 'Microsoft YaHei')

def create_risk_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, 0.5, 0.3, 9.0, 0.6, '偏差处理与风控', 28, ACCENT_RED, True, PP_ALIGN.CENTER, 'Microsoft YaHei')

    # Bias handling
    add_card(slide, 0.3, 1.2, 4.4, 3.0)
    add_text_box(slide, 0.5, 1.3, 4.0, 0.4, '偏差处理', 16, ACCENT_BLUE, True, font_name='Microsoft YaHei')

    biases = [
        ('幸存者偏差', 'D/G: 已消除 (baostock历史成分股)\nE: 中等 (手选龙头池, 2-4% CAGR影响)', ACCENT_GREEN),
        ('前视偏差', '全部消除 (因子仅用历史数据)', ACCENT_GREEN),
        ('交易成本', 'A/D/G: 8bps | F: 15bps | B/C: 5bps', ACCENT_GREEN),
        ('数据覆盖', 'D: 445只 | G: 928只 (PB覆盖496只)', TEXT_GRAY),
    ]
    for i, (title, desc, color) in enumerate(biases):
        y = 1.8 + i * 0.6
        add_text_box(slide, 0.5, y, 1.5, 0.25, title, 10, color, True, font_name='Microsoft YaHei')
        add_text_box(slide, 0.5, y+0.22, 4.0, 0.4, desc, 8, TEXT_GRAY, font_name='Microsoft YaHei')

    # Risk management
    add_card(slide, 5.0, 1.2, 4.7, 3.0)
    add_text_box(slide, 5.2, 1.3, 4.3, 0.4, '风控要点', 16, ACCENT_RED, True, font_name='Microsoft YaHei')

    risks = [
        '策略间低相关: 债/A股/港股/美股分散',
        'D+G: 低波策略天然控回撤 (MDD<12%)',
        'F: 双均线+回撤锁定自动切避险',
        'B: Vol缩放 + Regime过滤双重保护',
        'E: 仅买龙头 + 跌幅>20%过滤极端风险',
        '建议: 确认投入后逐步建仓, 非一次性满仓',
    ]
    for i, text in enumerate(risks):
        y = 1.8 + i * 0.38
        add_text_box(slide, 5.2, y, 4.3, 0.35, f'  {text}', 9, TEXT_LIGHT, font_name='Microsoft YaHei')

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    create_title_slide(prs)
    create_summary_slide(prs)
    create_strategy_a_slide(prs)
    create_strategy_b_slide(prs)
    create_strategy_c_slide(prs)
    create_strategy_d_slide(prs)
    create_strategy_e_slide(prs)
    create_strategy_f_slide(prs)
    create_strategy_g_slide(prs)
    create_risk_slide(prs)

    out_path = os.path.join(DATA_DIR, 'strategy_position_report_20260324.pptx')
    prs.save(out_path)
    print(f'PPT saved to: {out_path}')
    return out_path

if __name__ == '__main__':
    main()
